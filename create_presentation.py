"""Generate Voice360-Capstone.pptx from presentation-5-slides.md content."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colours and fonts ─────────────────────────────────────────────────────────
NAVY = RGBColor(0x0D, 0x1B, 0x2A)
BLUE = RGBColor(0x1A, 0x6F, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW = RGBColor(0xFF, 0xD7, 0x00)

TITLE_SIZE = 28
BODY_SIZE = 16
TABLE_SIZE = 13


def set_slide_background(slide, colour):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def add_title(slide, text, left=0.4, top=0.2, width=12.0, height=0.8):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(TITLE_SIZE)
    p.font.bold = True
    p.font.color.rgb = YELLOW


def add_textbox(slide, text, left, top, width, height, size=BODY_SIZE, bold=False, colour=WHITE):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = colour
    return txb


def add_bullet_list(slide, heading, items, left, top, width, height):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = heading
    p.font.name = "Calibri"
    p.font.size = Pt(BODY_SIZE)
    p.font.bold = True
    p.font.color.rgb = YELLOW
    for item in items:
        para = tf.add_paragraph()
        para.text = f"• {item}"
        para.font.name = "Calibri"
        para.font.size = Pt(BODY_SIZE)
        para.font.color.rgb = WHITE


def add_table(slide, headers, rows, left, top, width, height):
    cols = len(headers)
    row_count = len(rows) + 1
    tbl = slide.shapes.add_table(row_count, cols, Inches(left), Inches(top),
                                  Inches(width), Inches(height)).table
    tbl.columns[0].width = Inches(width * 0.38)
    if cols > 1:
        tbl.columns[1].width = Inches(width * 0.62)

    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.name = "Calibri"
        p.font.size = Pt(TABLE_SIZE)
        p.font.bold = True
        p.font.color.rgb = WHITE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid()
            fill_colour = RGBColor(0x0F, 0x24, 0x36) if ri % 2 == 0 else RGBColor(0x12, 0x2B, 0x42)
            cell.fill.fore_color.rgb = fill_colour
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.name = "Calibri"
            p.font.size = Pt(TABLE_SIZE)
            p.font.color.rgb = WHITE


# ── Build presentation ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ── SLIDE 1: Context and Why Now ──────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_background(slide, NAVY)
add_title(slide, "Slide 1: Context and Why Now")
add_textbox(slide, "Insurance servicing must move from queues to conversations",
            0.4, 0.9, 12.5, 0.5, size=20, bold=True)

add_bullet_list(slide, "Challenges and Impact", [
    "Customers wait in call queues for routine questions (claim status, missing docs)",
    "Rigid IVR menus cannot understand context or complete actions",
    "Contact-centre teams face rising demand, cost pressure, and high turnover",
    "Customers must repeat identity and context when transferred",
], 0.4, 1.55, 6.0, 2.4)

add_bullet_list(slide, "Strategic Imperative", [
    "Customers expect immediate, natural, always-available service",
    "Human capacity should focus on sensitive and complex cases",
    "Generative AI + controlled automation can combine conversation, knowledge, and action",
], 6.6, 1.55, 6.3, 1.8)

add_bullet_list(slide, "Shift in Operating Model", [
    "FROM: Queue-based, agent-only servicing",
    "TO: AI-assisted self-service with humans handling judgment and exceptions",
], 0.4, 4.1, 12.5, 1.2)

add_textbox(slide, "Suggested visual: Before/after split — IVR queue on left | conversational resolution on right",
            0.4, 5.4, 12.5, 0.5, size=11, colour=RGBColor(0x80, 0xA0, 0xC0))

add_textbox(slide, "1 / 5", 12.5, 7.0, 0.8, 0.4, size=11,
            colour=RGBColor(0x80, 0xA0, 0xC0))

# ── SLIDE 2: Current Challenges and Gaps ─────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_background(slide, NAVY)
add_title(slide, "Slide 2: Current Challenges and Gaps")
add_textbox(slide, "Routine claim enquiries consume time but deliver a fragmented experience",
            0.4, 0.9, 12.5, 0.5, size=18, bold=True)

add_bullet_list(slide, "Current State (6 Steps)", [
    "1. Customer calls insurer  [WAIT]",
    "2. Customer navigates IVR menu  [WAIT]",
    "3. Agent asks identity questions again",
    "4. Agent searches across multiple systems  [WAIT]",
    "5. Agent explains claim status verbally",
    "6. Customer repeats context to next agent  [WAIT if transferred]",
], 0.4, 1.55, 5.5, 3.6)

add_table(slide, ["Gap", "Description"], [
    ["High manual effort", "Repeatable enquiries consume agent time"],
    ["Long wait times", "No 24/7 self-service for routine queries"],
    ["Inconsistent explanations", "Claim stages explained differently by different agents"],
    ["Weak context transfer", "Customer must repeat everything when transferred"],
    ["Technology limitation", "Traditional automation can route but cannot safely retrieve individual status"],
], 6.2, 1.55, 6.7, 3.6)

add_textbox(slide, "2 / 5", 12.5, 7.0, 0.8, 0.4, size=11,
            colour=RGBColor(0x80, 0xA0, 0xC0))

# ── SLIDE 3: Vision and Future State ─────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_background(slide, NAVY)
add_title(slide, "Slide 3: Vision and Future State")
add_textbox(slide, "Voice360 resolves routine needs and brings people in at the right moment",
            0.4, 0.9, 12.5, 0.5, size=18, bold=True)

add_bullet_list(slide, "What the MVP Enables", [
    "Natural-language claim assistance — no menus, no navigation",
    "Approved, consistent answers to common process questions",
    "Secure retrieval of an existing claim after identity verification",
    "Clear missing-document and next-step guidance",
    "Callback creation with explicit customer confirmation",
    "Contextual human escalation — context travels with the customer",
], 0.4, 1.55, 5.8, 4.0)

add_table(slide, ["AI Handles", "People Handle"], [
    ["Intent recognition", "Coverage and liability decisions"],
    ["Approved FAQ responses", "Formal complaints and legal cases"],
    ["Identity verification flow", "Customer vulnerability and bereavement"],
    ["Claim data retrieval", "Fraud investigation"],
    ["Plain-language status explanation", "Disputed decisions"],
    ["Routine callback creation", "Any case requiring judgment"],
    ["Handoff summarisation", "Complex negotiations"],
], 6.4, 1.55, 6.5, 4.0)

add_textbox(slide, "3 / 5", 12.5, 7.0, 0.8, 0.4, size=11,
            colour=RGBColor(0x80, 0xA0, 0xC0))

# ── SLIDE 4: Solution Overview ────────────────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_background(slide, NAVY)
add_title(slide, "Slide 4: Solution Overview")
add_textbox(slide, "A grounded Copilot Studio agent with controlled actions and human guardrails",
            0.4, 0.9, 12.5, 0.5, size=18, bold=True)

arch = (
    "ARCHITECTURE\n"
    "Customer (Chat/Voice) → Voice360 — Microsoft Copilot Studio (Generative Orchestration)\n"
    "  ├─ Approved Claim Knowledge File (04-insurance-knowledge-base.md)\n"
    "  ├─ Power Automate Agent Flows: VerifyCustomer · GetClaimStatus · CreateCallbackRequest\n"
    "  │    └─ Voice360DemoData.xlsx (Customers / Claims / Callbacks — synthetic)\n"
    "  └─ Escalation Summary Flow → Human Agent or Callback\n"
    "GitHub Copilot ··· supports design, flow expressions, test generation, documentation"
)
add_textbox(slide, arch, 0.4, 1.55, 5.8, 2.8, size=11)

add_table(slide, ["Capability", "How it works"], [
    ["Generative orchestration", "Selects knowledge, topic, or action per turn"],
    ["Grounded knowledge", "Answers only from approved knowledge — no web browsing"],
    ["Deterministic verification", "All three identity values must match before any data is revealed"],
    ["Claim ownership check", "Flow filters by both customer ID and claim number"],
    ["Confirmation gate", "Explicit user confirmation required before callback is created"],
    ["Structured escalation", "PrepareHandoffSummary captures full context before transfer"],
], 6.4, 1.55, 6.5, 2.8)

add_table(slide, ["AI Tool", "Role in this project"], [
    ["Microsoft Copilot Studio", "Conversational agent, generative orchestration, topic and action hosting"],
    ["GitHub Copilot", "Requirements, expressions, test generation, security review, documentation"],
], 0.4, 4.55, 12.5, 1.4)

add_textbox(slide, "4 / 5", 12.5, 7.0, 0.8, 0.4, size=11,
            colour=RGBColor(0x80, 0xA0, 0xC0))

# ── SLIDE 5: Business Flow and Role of AI ────────────────────────────────────
slide = prs.slides.add_slide(blank_layout)
set_slide_background(slide, NAVY)
add_title(slide, "Slide 5: Business Flow and Role of AI")
add_textbox(slide, "One end-to-end journey from question to resolution or contextual handoff",
            0.4, 0.9, 12.5, 0.5, size=18, bold=True)

add_table(slide, ["Step", "What AI does"], [
    ["1", "Recognises customer intent — FAQ vs personal query"],
    ["2", "Retrieves answer from approved knowledge (not web)"],
    ["3", "Collects verification inputs conversationally"],
    ["4", "Calls VerifyCustomer action and handles result"],
    ["5", "Calls GetClaimStatus with verified customer ID only"],
    ["6", "Explains returned status in plain language"],
    ["7", "Prepares structured handoff summary before escalation"],
], 0.4, 1.55, 6.0, 3.2)

add_bullet_list(slide, "Control Points — What AI Does NOT Do", [
    "❌ No claim decisions — AI cannot approve, deny, or interpret liability",
    "❌ No data without verification — claim info only after verified=true",
    "❌ No cross-customer access — flow filters by customer ID in every query",
    "❌ No write actions without confirmation — callback requires confirmed=true",
    "❌ No sensitive data exposure — phone, DOB, full policy never echoed back",
], 6.6, 1.55, 6.3, 3.2)

add_textbox(slide,
            '"Voice360 makes routine motor claim servicing immediate and consistent,\n'
            'while keeping consequential insurance decisions firmly with people."',
            0.4, 5.0, 12.5, 1.2, size=15, bold=True)

add_textbox(slide, "5 / 5", 12.5, 7.0, 0.8, 0.4, size=11,
            colour=RGBColor(0x80, 0xA0, 0xC0))

# ── Save ──────────────────────────────────────────────────────────────────────
output = Path("Voice360-Capstone.pptx")
prs.save(output)
print(f"Created: {output}")
